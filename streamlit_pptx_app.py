# Streamlit app to generate and download a PowerPoint presentation for Cloud Security Transformation with CCSP
# Prerequisites: Install dependencies (pip install -r requirements.txt)
# Usage: Save as streamlit_pptx_app.py, run with `streamlit run streamlit_pptx_app.py`, open http://localhost:8501,
# click "Generate and Download PPTX" to create and download cloud_security_ccsp.pptx

import streamlit as st
import io
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.chart.data import ChartData
except ImportError:
    st.error("The 'python-pptx' library is not installed. Install it by running: `pip install python-pptx==0.6.23`")
    st.error("If using Python 3, try: `pip3 install python-pptx==0.6.23`")
    st.error("For virtual environments, activate the environment first.")
    st.error("For permission issues, try: `pip install python-pptx==0.6.23 --user`")
    st.stop()

# Define colors to match HTML
PRIMARY_BLUE = RGBColor(102, 126, 234)  # #667eea
ACCENT_RED = RGBColor(245, 87, 108)     # #f5576c
GREEN = RGBColor(40, 167, 69)           # #28a745
PURPLE = RGBColor(118, 75, 162)         # #764ba2
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(248, 249, 250)          # #f8f9fa
DARK_GRAY = RGBColor(44, 62, 80)        # #2c3e50

# Helper function to add a slide with title and content
def add_content_slide(prs, layout, title_text, content_text, title_color=PRIMARY_BLUE):
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = content_text
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.LEFT
    return slide

# Helper function to add a formatted text box
def add_text_box(slide, text, left, top, width, height, bg_color=WHITE, text_color=DARK_GRAY, font_size=16, is_mono=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER if is_mono else PP_ALIGN.LEFT
    if is_mono:
        p.font.name = 'Courier New'
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = bg_color
    txBox.line.color.rgb = PRIMARY_BLUE
    txBox.line.width = Pt(1)
    return txBox

# Function to generate the PowerPoint presentation
def generate_pptx():
    prs = Presentation()

    # Slide 1: Welcome & Introduction
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "Cloud Security Transformation with CCSP"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE

    subtitle = slide.placeholders[1]
    subtitle.text = ("Securing the Cloud Journey\n\n"
                     "Cloud adoption fuels innovation, but security risks like misconfigurations, "
                     "API vulnerabilities, and compliance gaps threaten enterprises. CCSP equips all "
                     "levels—beginners to advanced professionals—to build resilient, compliant multi-cloud systems.")
    subtitle.text_frame.paragraphs[0].font.size = Pt(14)
    subtitle.text_frame.paragraphs[0].font.color.rgb = WHITE
    subtitle.fill.solid()
    subtitle.fill.fore_color.rgb = ACCENT_RED

    stats = [("High", "Cost of breaches"), ("Critical", "Multi-cloud governance"), ("Strategic", "CCSP-driven resilience")]
    for i, (num, desc) in enumerate(stats):
        add_text_box(slide, f"{num}\n{desc}", Inches(2.5 + i*3), Inches(4), Inches(2), Inches(1.5), PURPLE, WHITE, 14)

    # Slide 2: Cloud vs. On-Premises Risks
    content_layout = prs.slide_layouts[1]
    slide = add_content_slide(prs, content_layout, "Cloud vs. On-Premises Risks",
        "Basic and Advanced Risk Analysis\n"
        "Cloud-Specific Risks:\n"
        "• Basic: Shared responsibility confusion (e.g., AWS S3 settings); multi-tenancy risks.\n"
        "• Advanced: Inconsistent multi-cloud governance; API vulnerabilities (e.g., Lambda).\n"
        "• Examples: 2018 S3 leak (basic); 2020 Twilio API breach (advanced).\n"
        "• Implications: Beginners—Learn roles. Intermediate—Use CSPM. Advanced—Standardize (Domain 1).\n\n"
        "On-Premises Risks:\n"
        "• Basic: Physical vulnerabilities; unpatched legacy systems.\n"
        "• Advanced: Limited DR redundancy; manual governance.\n"
        "• Examples: Equifax 2017 (basic); outdated servers (advanced).\n"
        "• Implications: Beginners—Patch. Intermediate—Automate backups. Advanced—Hybrid monitoring (Domain 5).\n\n"
        "Shared Risks:\n"
        "• Basic: Data breaches; insider threats.\n"
        "• Advanced: APTs; global compliance (GDPR, CCPA).\n"
        "• Examples: Target 2013 (basic); 2021 SolarWinds (advanced).\n"
        "• Implications: Beginners—Basic controls. Intermediate—SIEM. Advanced—Hybrid (Domains 1, 5).")
    add_text_box(slide, "[Cloud Risks]         [Shared Risks]         [On-Prem Risks]\n"
                        " | S3 Misconfig |----| Breaches, APTs |----| Legacy Systems |\n"
                        " | APIs, Shadow IT|----| Insider, Compliance|----| Physical, Manual |\n"
                        " | Domain 1, 4   |----| Domain 1, 5    |----| Domain 5        |",
                 Inches(2), Inches(5), Inches(6), Inches(1.5), GRAY, DARK_GRAY, 12, True)

    # Slide 3: Interactive Poll (Static)
    slide = add_content_slide(prs, content_layout, "Your Biggest Cloud Security Concern",
        "Prioritize Challenges for All Levels\n\n"
        "Misconfigured Services (Domain 1):\n"
        "• Impact: Exposed data due to improper settings.\n"
        "• Example: 2018 S3 bucket leak.\n"
        "• Implications: Beginners—Check defaults. Intermediate—AWS Config. Advanced—Automate audits.\n\n"
        "Lack of Visibility (Domain 5):\n"
        "• Impact: Delayed threat detection.\n"
        "• Example: Target 2013 breach.\n"
        "• Implications: Beginners—Learn monitoring. Intermediate—Splunk. Advanced—UEBA (Sentinel).\n\n"
        "Multi-Cloud Governance (Domain 1):\n"
        "• Impact: Inconsistent policies.\n"
        "• Example: 2019 Capital One IAM.\n"
        "• Implications: Beginners—Learn governance. Intermediate—AWS Organizations. Advanced—Standardize.\n\n"
        "API and Serverless Security (Domain 4):\n"
        "• Impact: Exposed endpoints.\n"
        "• Example: 2020 Twilio API breach.\n"
        "• Implications: Beginners—API basics. Intermediate—OAuth. Advanced—API Gateway.")

    # Slide 4: Key Security Challenges
    slide = add_content_slide(prs, content_layout, "Key Security Challenges in Cloud",
        "Basic and Advanced CCSP Mitigations\n\n"
        "1. Shared Responsibility Confusion:\n"
        "• Issue: Misunderstanding roles. Example: 2018 S3 leak.\n"
        "• Domains: 1 (Responsibility), 5 (Configuration).\n"
        "• Mitigation: Basic—Domain 1 matrix. Intermediate—AWS Config. Advanced—Azure Policy.\n\n"
        "2. Multi-Cloud Governance:\n"
        "• Issue: Inconsistent policies. Example: Capital One 2019.\n"
        "• Domains: 1, 5.\n"
        "• Mitigation: Basic—Learn governance. Intermediate—AWS Organizations. Advanced—ServiceNow GRC.\n\n"
        "3. Data Loss and Leakage:\n"
        "• Issue: Exposed data. Example: 2018 voter leak.\n"
        "• Domains: 2 (Encryption), 3 (Storage).\n"
        "• Mitigation: Basic—Encrypt. Intermediate—AWS Macie. Advanced—Vault.\n\n"
        "4. API and Serverless Security:\n"
        "• Issue: Exposed endpoints. Example: 2020 Twilio.\n"
        "• Domains: 2, 4.\n"
        "• Mitigation: Basic—OAuth. Intermediate—API Gateway. Advanced—Snyk.\n\n"
        "5. Visibility and Threat Detection:\n"
        "• Issue: Delayed detection. Example: 2021 SolarWinds.\n"
        "• Domains: 3, 5.\n"
        "• Mitigation: Basic—Monitoring. Intermediate—Splunk. Advanced—Cortex XSOAR.")

    # Slide 5: CCSP Domains Overview
    slide = add_content_slide(prs, content_layout, "CCSP Domains for Transformation",
        "Six Pillars for All Levels\n\n"
        "1. Cloud Concepts, Architecture & Design:\n"
        "• Principle: Align models with business needs.\n"
        "• Example: Netflix hybrid (advanced); IaaS/PaaS (basic).\n"
        "• Tools: Beginner—AWS Config. Advanced—AWS Organizations, Azure Policy.\n\n"
        "2. Cloud Data Security:\n"
        "• Principle: Protect data with encryption, DLP.\n"
        "• Example: Dropbox AES-256 (basic); Vault (advanced).\n"
        "• Tools: AWS KMS, Vault.\n\n"
        "3. Cloud Platform & Infrastructure Security:\n"
        "• Principle: Secure networks, containers.\n"
        "• Example: VPC (basic); Kubernetes with Istio (advanced).\n"
        "• Tools: AWS VPC, Istio.\n\n"
        "4. Cloud Application Security:\n"
        "• Principle: Embed security in apps.\n"
        "• Example: Secure coding (basic); Snyk with GitHub Actions (advanced).\n"
        "• Tools: Code reviews, Snyk.\n\n"
        "5. Cloud Security Operations:\n"
        "• Principle: Automate monitoring, response.\n"
        "• Example: Splunk (basic); Cortex XSOAR (advanced).\n"
        "• Tools: Splunk, XSOAR.\n\n"
        "6. Legal, Risk & Compliance:\n"
        "• Principle: Ensure compliance.\n"
        "• Example: GDPR checks (basic); OneTrust (advanced).\n"
        "• Tools: Manual audits, OneTrust.")

    # Slide 6: Best Practices
    slide = add_content_slide(prs, content_layout, "Best Practices for Secure Cloud Adoption",
        "Phased Approach for All Levels\n\n"
        "Phase 1: Assessment (Domain 1, 6):\n"
        "• Action: Inventory assets; establish multi-cloud governance.\n"
        "• Tools: Beginner—AWS Config. Advanced—ServiceNow GRC, AWS Organizations.\n"
        "• Example: Retailer used Config; bank standardized policies.\n\n"
        "Phase 2: Implementation (Domain 2, 3, 4):\n"
        "• Action: Deploy encryption, RBAC, DevSecOps.\n"
        "• Tools: Beginner—AWS KMS, Okta. Advanced—Vault, Istio, Snyk.\n"
        "• Example: Tech firm used Okta; fintech secured serverless.\n\n"
        "Phase 3: Operations (Domain 5):\n"
        "• Action: Automate monitoring, incident response.\n"
        "• Tools: Beginner—Splunk. Advanced—Cortex XSOAR, Sentinel.\n"
        "• Example: Healthcare used Splunk; enterprise reduced MTTD by 80%.\n")
    add_text_box(slide, "Key Tools: CASB, CSPM, SOAR, GRC\nGoal: Compliant, Secure Operations",
                 Inches(2), Inches(5), Inches(6), Inches(0.8), GREEN, WHITE, 16)

    # Slide 7: Detailed Case Study with Chart
    slide = add_content_slide(prs, content_layout, "Case Study: Multi-Cloud Financial Transformation",
        "Securing 300+ Apps Across AWS, Azure, GCP\n\n"
        "Context: Global financial firm migrated 300+ apps to multi-cloud under SOX, PCI DSS, GDPR, CCPA.\n\n"
        "Stage 1: Planning (Domains 1, 6):\n"
        "• Requirements: 1M+ transactions; 100% visibility.\n"
        "• Risks: Basic—Shared responsibility (e.g., 2018 S3). Advanced—Governance (e.g., 2019 Capital One).\n"
        "• Security: Basic—Domain 1 matrix. Intermediate—AWS Config. Advanced—AWS Organizations, ServiceNow GRC.\n"
        "• Outcome: 100% visibility; GDPR readiness.\n"
        "• Tools: AWS Config, Organizations, ServiceNow GRC.\n\n"
        "Stage 2: Implementation (Domains 2, 3, 4):\n"
        "• Requirements: Secure 1PB+ data, 200+ serverless apps; zero-trust.\n"
        "• Risks: Basic—S3 misconfigs (e.g., 2018 voter). Advanced—API gaps (e.g., 2020 Twilio).\n"
        "• Security: Basic—KMS, VPCs. Intermediate—Macie, API Gateway. Advanced—Vault, Istio, Snyk.\n"
        "• Outcome: 90% attack surface reduction; zero-trust apps.\n"
        "• Tools: KMS, Macie, API Gateway, Vault, Istio, Snyk.\n\n"
        "Stage 3: Optimization (Domain 5):\n"
        "• Requirements: MTTD <24h, MTTR <12h; automate compliance.\n"
        "• Risks: Basic—Visibility (e.g., Target 2013). Advanced—APTs (e.g., 2021 SolarWinds).\n"
        "• Security: Basic—Splunk. Intermediate—Sentinel. Advanced—XSOAR.\n"
        "• Outcome: MTTD <24h; 75% MTTR improvement.\n"
        "• Tools: Splunk, Sentinel, XSOAR.\n\n"
        "Results (All Domains):\n"
        "• Impact: Resilient, compliant multi-cloud architecture.")

    # Add chart
    chart_data = ChartData()
    chart_data.categories = ['Planning', 'Implementation', 'Optimization', 'Overall']
    chart_data.add_series('Compliance (%)', (100, 90, 95, 100))
    chart_data.add_series('Attack Surface Reduction (%)', (0, 90, 80, 90))
    chart_data.add_series('MTTD Reduction (%)', (0, 0, 80, 80))
    chart_data.add_series('MTTR Improvement (%)', (0, 0, 75, 75))

    chart_placeholder = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(2), Inches(5), Inches(6), Inches(2.5), chart_data)
    chart = chart_placeholder.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Case Study Metrics by Stage"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(16)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.category_axis.has_title = True
    chart.category_axis.axis_title.text_frame.text = "Stages"
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = "Percentage (%)"
    chart.value_axis.maximum_scale = 100

    # Apply fill colors at series level
    colors = [PRIMARY_BLUE, ACCENT_RED, GREEN, PURPLE]
    for i, series in enumerate(chart.series):
        try:
            series.fill.solid()
            series.fill.fore_color.rgb = colors[i]
        except AttributeError:
            st.warning(f"Could not set fill for series {series.name}. Using default styling.")

    add_text_box(slide, "[Multi-Cloud Architecture]\n"
                        " | AWS (Config, KMS) |----| Azure (Sentinel, Policy) |----| GCP (RBAC, GRC) |\n"
                        " | Stage 1: Governance|----| Stage 2: Zero-Trust |----| Stage 3: SOAR     |\n"
                        " | Vault, Istio, Snyk |----| XSOAR, Compliance   |----| API Gateway       |\n"
                        " | Domain 1,2,3,4,5,6 |",
                 Inches(2), Inches(7.5), Inches(6), Inches(1.2), GRAY, DARK_GRAY, 10, True)

    # Slide 8: Interactive Quiz (Static)
    slide = add_content_slide(prs, content_layout, "Test Your CCSP Knowledge",
        "Question 1: Who secures data in AWS S3? (Domain 1)\n"
        "• Options: AWS Only, Customer, Both Equally\n"
        "• Correct: Customer. The customer secures S3 data.\n\n"
        "Question 2: Which tool supports Domain 5’s incident response?\n"
        "• Options: Palo Alto Cortex XSOAR, AWS KMS, OneTrust\n"
        "• Correct: Cortex XSOAR aligns with Domain 5 for automation.")

    # Slide 9: Actionable Steps
    slide = add_content_slide(prs, content_layout, "Actionable Steps for Cloud Security",
        "Implement CCSP-Driven Security Today\n\n"
        "Beginner: Enable MFA (Domain 3):\n"
        "• Activate MFA on AWS IAM or Azure AD.\n"
        "• Action: Enable MFA for admin accounts.\n"
        "• Tool: AWS IAM, Okta.\n\n"
        "Beginner: Secure S3 Buckets (Domain 2):\n"
        "• Restrict public access to prevent leaks.\n"
        "• Action: Review S3 permissions.\n"
        "• Tool: AWS S3 Console.\n\n"
        "Intermediate: Deploy CSPM (Domain 5):\n"
        "• Monitor and fix misconfigurations.\n"
        "• Action: Set up AWS Config or Azure Security Center.\n"
        "• Tool: AWS Config, Azure Security Center.\n\n"
        "Intermediate: Monitor with SIEM (Domain 5):\n"
        "• Real-time threat detection.\n"
        "• Action: Configure Splunk for logs and alerts.\n"
        "• Tool: Splunk, Azure Sentinel.\n\n"
        "Advanced: Standardize Governance (Domain 1):\n"
        "• Consistent policies across clouds.\n"
        "• Action: Use AWS Organizations, ServiceNow GRC.\n"
        "• Tool: AWS Organizations, ServiceNow GRC.\n\n"
        "Advanced: Automate Incident Response (Domain 5):\n"
        "• Deploy SOAR for threat response.\n"
        "• Action: Integrate Cortex XSOAR with Sentinel.\n"
        "• Tool: Cortex XSOAR, Palo Alto Sentinel.")

    # Slide 10: Engage and Learn More
    slide = add_content_slide(prs, content_layout, "Strategize and Scale with CCSP",
        "Next Steps for Cloud Security\n\n"
        "Domain 1:\n"
        "• Assess cloud models and governance (AWS Organizations).\n\n"
        "Domain 2–4:\n"
        "• Implement encryption (KMS, Vault), RBAC (Istio), DevSecOps (Snyk).\n\n"
        "Domain 5–6:\n"
        "• Automate monitoring (XSOAR, Splunk) and compliance (OneTrust).\n\n"
        "Action: Enroll in CCSP training at www.isc2.org.\n\n"
        "Enhanced: Resilience with CCSP\n"
        "Unified: Multi-cloud strategy")
    add_text_box(slide, "Enhanced: Resilience with CCSP\nUnified: Multi-cloud strategy",
                 Inches(2), Inches(7), Inches(6), Inches(0.8), PURPLE, WHITE, 16)

    # Save to a temporary file
    pptx_file = "cloud_security_ccsp.pptx"
    prs.save(pptx_file)
    return pptx_file

# Streamlit app
st.title("Cloud Security CCSP Presentation Generator")
st.write("Click the button below to generate and download a PowerPoint presentation on Cloud Security Transformation with CCSP.")

if st.button("Generate and Download PPTX"):
    with st.spinner("Generating PowerPoint file..."):
        pptx_file = generate_pptx()
        with open(pptx_file, "rb") as f:
            pptx_bytes = f.read()
        st.download_button(
            label="Download cloud_security_ccsp.pptx",
            data=pptx_bytes,
            file_name="cloud_security_ccsp.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        st.success("Presentation generated! Click the download button above.")
        # Clean up temporary file
        if os.path.exists(pptx_file):
            os.remove(pptx_file)
